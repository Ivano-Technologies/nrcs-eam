#!/usr/bin/env python3
"""Build 60-slide NRCS SEIP investor PowerPoint from generated PNG assets."""
from pathlib import Path
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

OUTPUT_DIR = Path("C:/mnt/agents/output")
if not OUTPUT_DIR.exists():
    OUTPUT_DIR = Path("/mnt/agents/output")
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

OUTPUT_FILE = OUTPUT_DIR / "NRCS_SEIP_Investor_Presentation.pptx"

# Design tokens
DEEP_BLUE = RGBColor(0x00, 0x28, 0x55)
RED = RGBColor(0xE3, 0x18, 0x37)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CHARCOAL = RGBColor(0x36, 0x45, 0x4F)
BLACK = RGBColor(0x00, 0x00, 0x00)
SLATE = RGBColor(0x70, 0x80, 0x90)
LIGHT_BLUE = RGBColor(0xE6, 0xF3, 0xFF)
FOOTER_TEXT = "STRICTLY PRIVATE AND CONFIDENTIAL | NRCS SEIP | May 2026"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

ASSET_MAP = {
    "hero": "01_hero_ambulance_deficit.png",
    "title": "02_title_slide.png",
    "flywheel": "03_flywheel.png",
    "market": "04_market_pyramid.png",
    "competitive": "05_competitive_matrix.png",
    "financial_dash": "06_financial_dashboard.png",
    "risk": "07_risk_heat_map.png",
    "gantt": "08_gantt.png",
    "sdg": "09_sdg_dashboard.png",
    "revenue": "10_revenue_deepdive.png",
    "governance": "11_governance.png",
    "sensitivity": "12_sensitivity.png",
    "leadership": "13_leadership.png",
    "donor": "14_donor_crisis.png",
    "features": "15_feature_matrix.png",
    "exec": "16_exec_summary.png",
    "ecosystem": "17_ecosystem.png",
    "ambulance": "18_ambulance_model.png",
    "assets": "19_asset_mgmt.png",
    "fundraising": "20_fundraising.png",
    "nti": "21_nti_courses.png",
    "gtm": "22_gtm_strategy.png",
    "hospitality": "23_hospitality_gap.png",
    "statements": "24_financial_statements.png",
}


def asset_path(key):
    return OUTPUT_DIR / ASSET_MAP[key]


def missing_assets():
    missing = []
    for key, fname in ASSET_MAP.items():
        p = OUTPUT_DIR / fname
        if not p.exists():
            missing.append(fname)
    return missing


def set_slide_size(prs):
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H


def add_header_bar(slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.35))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DEEP_BLUE
    shape.line.fill.background()
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.35), SLIDE_W, Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RED
    accent.line.fill.background()


def add_footer(slide, slide_num, total=60):
    box = slide.shapes.add_textbox(Inches(0.3), Inches(7.05), Inches(10), Inches(0.35))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = FOOTER_TEXT
    p.font.size = Pt(8)
    p.font.name = "Arial"
    p.font.color.rgb = SLATE
    p.alignment = PP_ALIGN.RIGHT
    num = slide.shapes.add_textbox(Inches(12.3), Inches(7.05), Inches(0.8), Inches(0.35))
    np = num.text_frame.paragraphs[0]
    np.text = f"{slide_num}/{total}"
    np.font.size = Pt(9)
    np.font.name = "Arial"
    np.font.color.rgb = SLATE
    np.alignment = PP_ALIGN.RIGHT


def add_title(slide, text, top=Inches(0.55)):
    box = slide.shapes.add_textbox(Inches(0.5), top, Inches(12), Inches(0.7))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.name = "Arial"
    p.font.color.rgb = DEEP_BLUE


def add_body(slide, text, left=Inches(0.5), top=Inches(1.4), width=Inches(12), size=12):
    box = slide.shapes.add_textbox(left, top, width, Inches(5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.name = "Arial"
    p.font.color.rgb = CHARCOAL
    return box


def add_image(slide, key, left=0, top=0, width=None, height=None):
    path = asset_path(key)
    if not path.exists():
        add_body(slide, f"[Missing: {path.name}]", top=Inches(2))
        return False
    w = width or SLIDE_W
    h = height or SLIDE_H
    slide.shapes.add_picture(str(path), left, top, width=w, height=h)
    return True


def add_notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.text = text


def style_table(table, header_cols):
    for c in range(header_cols):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = DEEP_BLUE
        for p in cell.text_frame.paragraphs:
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.font.size = Pt(10)
            p.font.name = "Arial"
    for r in range(1, len(table.rows)):
        for c in range(header_cols):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BLUE if r % 2 == 0 else WHITE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.name = "Consolas" if c > 0 else "Arial"
                p.font.color.rgb = BLACK


def blank_slide(prs):
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)


def slide_full_bleed(prs, key, notes, num):
    slide = blank_slide(prs)
    add_image(slide, key, 0, 0, SLIDE_W, SLIDE_H)
    add_footer(slide, num)
    add_notes(slide, notes)
    return slide


def slide_with_image(prs, title, key, body, notes, num, img_top=Inches(1.3), img_h=Inches(4.8)):
    slide = blank_slide(prs)
    add_header_bar(slide)
    add_title(slide, title)
    add_image(slide, key, Inches(0.4), img_top, Inches(8.5), img_h)
    add_body(slide, body, left=Inches(9.2), top=Inches(1.4), width=Inches(3.8), size=11)
    add_footer(slide, num)
    add_notes(slide, notes)
    return slide


def slide_text(prs, title, body, notes, num):
    slide = blank_slide(prs)
    add_header_bar(slide)
    add_title(slide, title)
    add_body(slide, body, top=Inches(1.5), size=13)
    add_footer(slide, num)
    add_notes(slide, notes)
    return slide


def slide_nsia_table(prs, num):
    slide = blank_slide(prs)
    add_header_bar(slide)
    add_title(slide, "The NSIA Ask — CHF 750,000")
    rows, cols = 6, 3
    tbl_shape = slide.shapes.add_table(rows, cols, Inches(1), Inches(1.8), Inches(11), Inches(4.5))
    table = tbl_shape.table
    data = [
        ["Allocation", "Amount (CHF)", "% of Total"],
        ["Ambulance Fleet & Equipment", "225,000", "30%"],
        ["Gwarimpa & Utako Facilities", "262,500", "35%"],
        ["NTI Training Centre", "150,000", "20%"],
        ["CRM & Fundraising Platform", "75,000", "10%"],
        ["Working Capital & Contingency", "37,500", "5%"],
    ]
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            table.cell(r, c).text = val
    style_table(table, 3)
    add_body(slide, "5-year co-investment with NSIA oversight. Mission dividend: 12% of net surplus.", top=Inches(6.2), size=11)
    add_footer(slide, num)
    add_notes(slide, "Walk through allocation table. Emphasize 79% grant utilization and capex-heavy Year 1 deployment.")
    return slide


def slide_platform_quadrant(prs, num):
    slide = blank_slide(prs)
    add_header_bar(slide)
    add_title(slide, "Platform Overview — Four Verticals")
    quads = [
        ("NTI Training", "25% grant", "CHF 26–335/course\nMDCN accreditation path", Inches(0.5), Inches(1.5)),
        ("Ambulance Services", "30% grant", "PAG + Subscription\n10→25 fleet units", Inches(6.8), Inches(1.5)),
        ("Asset Management", "35% grant", "Gwarimpa halls + Utako Lodge\n35%→85% utilisation", Inches(0.5), Inches(4.2)),
        ("Fundraising CRM", "10% grant", "Salesforce + Paystack\nBronze→Platinum tiers", Inches(6.8), Inches(4.2)),
    ]
    for title, grant, detail, left, top in quads:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.8), Inches(2.4))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BLUE
        box.line.color.rgb = DEEP_BLUE
        tf = box.text_frame
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = DEEP_BLUE
        p2 = tf.add_paragraph()
        p2.text = grant
        p2.font.size = Pt(11)
        p2.font.color.rgb = RED
        p3 = tf.add_paragraph()
        p3.text = detail
        p3.font.size = Pt(10)
        p3.font.color.rgb = CHARCOAL
    add_footer(slide, num)
    add_notes(slide, "Introduce four integrated verticals. Each reinforces the flywheel and diversifies revenue away from donor dependency.")
    return slide


def financial_table(slide, title, headers, rows, left=Inches(0.5), top=Inches(1.5)):
    add_title(slide, title, top=Inches(0.55))
    r_count = len(rows) + 1
    c_count = len(headers)
    tbl = slide.shapes.add_table(r_count, c_count, left, top, Inches(12), Inches(0.35 * r_count)).table
    for c, h in enumerate(headers):
        tbl.cell(0, c).text = h
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            tbl.cell(ri + 1, ci).text = str(val)
    style_table(tbl, c_count)
    return tbl


def build_presentation():
    prs = Presentation()
    set_slide_size(prs)
    n = 0

    # SECTION 0 — OPENING (1-3)
    configs = [
        ("title", "Opening title. Set institutional tone. Introduce NRCS SEIP as NSIA strategic partnership opportunity."),
        ("hero", "Lead with 96% ambulance deficit statistic. Pause on WHO 1.0 vs Nigeria 0.4 comparison."),
        ("donor", "Frame donor dependency crisis. Transition to self-sustaining social enterprise solution."),
    ]
    for key, note in configs:
        n += 1
        slide_full_bleed(prs, key, note, n)

    # SECTION 1 — EXECUTIVE SUMMARY (4-8)
    n += 1
    slide_with_image(prs, "Executive Summary",
        "exec",
        "Integrated humanitarian infrastructure platform.\n4 verticals | CHF 750K ask | 79% grant utilisation\nBreak-even Month 18 | 12% mission dividend",
        "Summarise investment thesis in 90 seconds. Hit four metrics before deep dive.", n)

    n += 1
    slide_platform_quadrant(prs, n)

    n += 1
    slide_with_image(prs, "Flywheel Architecture", "flywheel",
        "Six interconnected nodes drive network effects:\nAmbulance → Training → Events → Lodge → Donors → Relationships",
        "Explain virtuous cycle. Each vertical feeds the next.", n)

    n += 1
    slide_with_image(prs, "Key Financial Highlights", "financial_dash",
        "Y5 Revenue: CHF 1.1M\nY5 EBITDA Margin: 26%\nBreak-even: Month 18\n5-year cumulative EBITDA: CHF 545K",
        "Walk through four dashboard panels. Emphasise margin inflection at Year 3.", n)

    n += 1
    slide_nsia_table(prs, n)

    # SECTION 2 — INDUSTRY OVERVIEW (9-13)
    industry = [
        ("Market Opportunity", "market", "TAM NGN 957B–2.12TN. SAM NGN 74.5B–188.2B. SOM NGN 1.3B–3.18B."),
        ("Hospitality Gap — Abuja", "hospitality", "Missing middle segment: NGOs, corporates, government delegations underserved."),
        ("Healthcare Infrastructure Gap", "hero", "Emergency response capacity at 7% population coverage."),
        ("Social Enterprise Landscape", "ecosystem", "Six stakeholder clusters enable platform scale."),
        ("Industry Growth Drivers", "market", "Healthcare gap, NGO demand, corporate CSR, government outsourcing."),
    ]
    for title, key, body in industry:
        n += 1
        slide_with_image(prs, title, key, body, f"Cover {title}. Use data points from visual.", n)

    # SECTION 3 — MARKET & COMPETITION (14-18)
    market_slides = [
        ("Competitive Positioning", "competitive", "NRCS SEIP occupies high-speed, high-integration quadrant."),
        ("Feature Comparison", "features", "Only NRCS SEIP delivers all 5 capability pillars."),
        ("Stakeholder Ecosystem", "ecosystem", "Government, international, corporate, community, hospitality, technology."),
        ("Competitive Moat", "flywheel", "Integrated flywheel creates switching costs and network effects."),
        ("Market Entry Strategy", "gtm", "First-mover in missing-middle hospitality and integrated EMS."),
    ]
    for title, key, body in market_slides:
        n += 1
        slide_with_image(prs, title, key, body, f"Present {title} with competitive narrative.", n)

    # SECTION 4 — PLATFORM DETAIL (19-25)
    platform = [
        ("Ambulance Vertical", "ambulance", "PAG + Subscription model. FAAN, FRSC, oil & gas, UN targets."),
        ("NTI Training Institute", "nti", "5 courses, CHF 26–335 pricing ladder. MDCN/ILCOR accreditation."),
        ("Asset Management", "assets", "Gwarimpa halls + Utako Lodge 16 rooms. 35%→85% utilisation."),
        ("Fundraising Platform", "fundraising", "Salesforce CRM. Bronze–Platinum corporate tiers."),
        ("Revenue Deep Dive", "revenue", "Waterfall, unit economics, pricing ladder, occupancy ramp."),
        ("Go-To-Market Strategy", "gtm", "17 channels across 5 verticals."),
        ("Platform Integration", "flywheel", "Cross-vertical synergies drive LTV/CAC improvement."),
    ]
    for title, key, body in platform:
        n += 1
        slide_with_image(prs, title, key, body, f"Deep dive: {title}.", n)

    # SECTION 5 — BUSINESS MODEL (26-30)
    biz_slides = [
        ("Revenue Streams", "Four vertical revenue contribution at Year 5."),
        ("Pricing Model", "NTI ladder CHF 26–335. Lodge ADR NGN 45K–65K."),
        ("Customer Acquisition", "17 GTM channels. Anchor contract strategy."),
        ("Partnership Ecosystem", "NSIA oversight + government/international partners."),
        ("Unit Economics", "LTV/CAC 0.8x → 3.1x. Payback < 14 months by Y3."),
    ]
    biz_keys = ["revenue", "nti", "gtm", "ecosystem", "revenue"]
    for (title, body), key in zip(biz_slides, biz_keys):
        n += 1
        slide_with_image(prs, title, key, body, f"Business model: {title}.", n)

    # SECTION 6 — MANAGEMENT (31-35)
    mgmt = [
        ("Leadership Team", "leadership", "6 executives with humanitarian, medical, ops, tech, finance, partnerships expertise."),
        ("Governance Structure", "governance", "NECS → SE Board → SEIP MD → 4 GMs. NSIA oversight."),
        ("Board Composition", "governance", "5-member SE Board with independent audit committee."),
        ("Advisory Network", "leadership", "MDCN, ILCOR, NCDMB, FMOH accreditation partners."),
        ("Organisational Culture", "sdg", "Mission-driven with 12% community reinvestment dividend."),
    ]
    for title, key, body in mgmt:
        n += 1
        slide_with_image(prs, title, key, body, f"Management: {title}.", n)

    # SECTION 7 — OPERATIONS (36-40)
    ops = [
        ("Operational Strategy", "gantt", "24-month phased implementation."),
        ("Facilities Plan", "assets", "Gwarimpa fit-out + Utako renovation."),
        ("Workforce Plan", "nti", "320 jobs by Year 5. NTI as talent pipeline."),
        ("Quality Assurance", "risk", "ISO-aligned protocols. 15% indigent quota."),
        ("Supply Chain", "ambulance", "Fleet SLA contracts. Maintenance partnerships."),
    ]
    for title, key, body in ops:
        n += 1
        slide_with_image(prs, title, key, body, f"Operations: {title}.", n)

    # SECTION 8 — FINANCIAL PLAN (41-48)
    n += 1
    slide_with_image(prs, "Financial Statements Overview", "statements",
        "Consolidated 5-year view: Revenue CHF 120K→1.1M. EBITDA -40K→290K.",
        "Overview of five financial panels before detailed tables.", n)

    n += 1
    slide_with_image(prs, "Sensitivity Analysis", "sensitivity",
        "Base case CHF 1.1M revenue. Upside CHF 1.45M. Occupancy is #1 sensitivity driver.",
        "Present tornado chart and three scenarios.", n)

    fin_tables = [
        ("5-Year P&L (CHF '000)", ["", "Y1", "Y2", "Y3", "Y4", "Y5"],
         [["Revenue", "120", "280", "520", "780", "1100"],
          ["Grant Income", "80", "60", "50", "40", "30"],
          ["Gross Profit", "60", "168", "338", "546", "814"],
          ["EBITDA", "-40", "20", "95", "180", "290"],
          ["PAT", "-55", "-10", "45", "110", "195"]]),
        ("Balance Sheet (CHF '000)", ["", "Y1", "Y3", "Y5"],
         [["Total Assets", "720", "980", "1250"],
          ["Fixed Assets", "580", "750", "820"],
          ["Current Assets", "140", "230", "430"],
          ["Total Liabilities", "180", "220", "280"],
          ["Equity", "540", "760", "970"]]),
        ("Cash Flow (CHF '000)", ["", "Y1", "Y2", "Y3", "Y4", "Y5"],
         [["Operating CF", "-30", "35", "88", "155", "220"],
          ["Investing CF", "-637", "-80", "-60", "-50", "-40"],
          ["Financing CF", "750", "0", "0", "0", "0"],
          ["Closing Cash", "83", "38", "66", "171", "351"]]),
        ("Capex Allocation", ["Category", "CHF '000", "%"],
         [["Ambulance Fleet", "225", "35%"],
          ["Facilities", "263", "41%"],
          ["NTI Centre", "90", "14%"],
          ["CRM Platform", "37", "6%"],
          ["Contingency", "22", "4%"]]),
        ("Key Metrics Summary", ["Metric", "Y1", "Y3", "Y5"],
         [["EBITDA Margin %", "-33", "18", "26"],
          ["Gross Margin %", "50", "65", "74"],
          ["Revenue Growth %", "—", "86", "41"],
          ["ROE %", "N/A", "6", "28"]]),
        ("Revenue Mix by Vertical (CHF '000)", ["Vertical", "Y1", "Y3", "Y5"],
         [["NTI", "30", "60", "80"],
          ["Ambulance", "40", "150", "280"],
          ["Asset Mgmt", "35", "220", "480"],
          ["Fundraising", "15", "90", "260"]]),
    ]
    for title, headers, rows in fin_tables:
        n += 1
        slide = blank_slide(prs)
        add_header_bar(slide)
        financial_table(slide, title, headers, rows)
        add_footer(slide, n)
        add_notes(slide, f"Review {title}. Highlight trajectory and key inflection points.")

    # SECTION 9 — RISK (49-52)
    risks = [
        ("Risk Register Overview", "risk", "8 risks rated 1–5. Donor fatigue and FX highest priority."),
        ("Regulatory & Compliance", "risk", "Licensing delays mitigated by early FMOH/MDCN engagement."),
        ("Operational Risks", "risk", "Fleet maintenance via SLA. 15% indigent quota managed."),
        ("Financial Risks", "sensitivity", "FX hedging. Scenario planning. Break-even M18."),
    ]
    for title, key, body in risks:
        n += 1
        slide_with_image(prs, title, key, body, f"Risk: {title}.", n)

    # SECTION 10 — IMPACT (53-55)
    impact = [
        ("SDG Impact Dashboard", "sdg", "SDG 3, 8, 9, 17 alignment with Year 1 vs Year 5 metrics."),
        ("Mission Dividend Framework", "sdg", "12% of net surplus reinvested in community programs."),
        ("Social Return on Investment", "sdg", "85K lives touched. 320 jobs. 45 partnerships by Y5."),
    ]
    for title, key, body in impact:
        n += 1
        slide_with_image(prs, title, key, body, f"Impact: {title}.", n)

    # SECTION 11 — APPENDIX (56-60)
    n += 1
    slide_text(prs, "Appendix: Market Sizing Methodology",
        "TAM: National healthcare + hospitality + training spend (NGN 957B–2.12TN).\n"
        "SAM: Abuja-FCT addressable market (NGN 74.5B–188.2B).\n"
        "SOM: NRCS SEIP realistic capture (NGN 1.3B–3.18B).\n"
        "Sources: WHO, NBS, industry interviews, comparable transactions.",
        "Reference for detailed market sizing questions.", n)

    n += 1
    slide = blank_slide(prs)
    add_header_bar(slide)
    financial_table(slide, "Appendix: Complete Financial Statements (CHF '000)",
        ["", "Y1", "Y2", "Y3", "Y4", "Y5"],
        [["Revenue", "120", "280", "520", "780", "1100"],
         ["COGS", "140", "172", "232", "274", "316"],
         ["Opex", "20", "88", "193", "326", "494"],
         ["EBITDA", "-40", "20", "95", "180", "290"],
         ["Net Income", "-55", "-10", "45", "110", "195"]],
        top=Inches(1.5))
    add_footer(slide, n)
    add_notes(slide, "Full P&L appendix for IC review.")

    n += 1
    slide = blank_slide(prs)
    add_header_bar(slide)
    add_title(slide, "Appendix: Assumptions Register")
    assumptions = [
        ["Assumption", "Value", "Source"],
        ["Ambulance utilisation Y5", "72%", "Anchor contract pipeline"],
        ["Lodge occupancy Y5", "82%", "Comparable Abuja lodges"],
        ["FX rate", "NGN 1,450/USD", "Conservative forward curve"],
        ["Discount rate", "12%", "Social enterprise benchmark"],
        ["Mission dividend", "12%", "NRCS board resolution"],
    ]
    tbl = slide.shapes.add_table(len(assumptions), 3, Inches(0.5), Inches(1.5), Inches(12), Inches(4)).table
    for r, row in enumerate(assumptions):
        for c, val in enumerate(row):
            tbl.cell(r, c).text = val
    style_table(tbl, 3)
    add_footer(slide, n)
    add_notes(slide, "Key assumptions for sensitivity and scenario analysis.")

    n += 1
    slide_text(prs, "Appendix: Glossary",
        "SEIP — Social Enterprise Infrastructure Platform\n"
        "NTI — National Training Institute\n"
        "PAG — Pay-As-Go ambulance billing\n"
        "NEC — National Executive Council\n"
        "NSIA — Nigeria Sovereign Investment Authority\n"
        "SAM/SOM/TAM — Serviceable/Obtainable/Total Addressable Market\n"
        "MDCN — Medical and Dental Council of Nigeria",
        "Quick reference for institutional reviewers.", n)

    n += 1
    slide_full_bleed(prs, "title",
        "Closing: Thank you. Q&A. Reiterate CHF 750K ask and 12% mission dividend commitment.", n)

    prs.save(str(OUTPUT_FILE))
    return n


if __name__ == "__main__":
    missing = missing_assets()
    if missing:
        print(f"WARNING: Missing assets before build: {missing}")
    total = build_presentation()
    size_mb = OUTPUT_FILE.stat().st_size / (1024 * 1024)
    incorporated = [f for f in ASSET_MAP.values() if (OUTPUT_DIR / f).exists()]
    print(f"\nPresentation saved: {OUTPUT_FILE}")
    print(f"File size: {size_mb:.2f} MB")
    print(f"Total slides: {total}")
    print(f"PNG assets incorporated: {len(incorporated)}/24")
    if missing:
        print(f"Missing assets: {missing}")
    else:
        print("All 24 PNG assets present and referenced.")
    errors = []
    from pptx import Presentation as P
    verify = P(str(OUTPUT_FILE))
    print(f"Verified slide count in file: {len(verify.slides)}")
    if len(verify.slides) != 60:
        errors.append(f"Expected 60 slides, got {len(verify.slides)}")
    if errors:
        print("Errors:", errors)
    else:
        print("Build completed successfully.")
